# ----------------------------------------------------------------------
# Wordless: File Area
# Copyright (C) 2018-2025  Ye Lei (叶磊)
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU General Public License for more details.
#
# You should have received a copy of the GNU General Public License
# along with this program.  If not, see <https://www.gnu.org/licenses/>.
# ----------------------------------------------------------------------

# pylint: disable=broad-exception-caught

import copy
import csv
import os
import re
import traceback

import bs4
import docx
import openpyxl
import pptx
import pypdf
from PyQt5.QtCore import (
    pyqtSignal,
    QCoreApplication,
    QItemSelection,
    QRect,
    Qt
)
from PyQt5.QtGui import QStandardItem
from PyQt5.QtWidgets import (
    QAbstractItemDelegate,
    QCheckBox,
    QFileDialog,
    QHeaderView,
    QLineEdit,
    QPushButton,
    QStyle,
    QStyleOptionButton
)

from wordless.wl_checks import wl_checks_files, wl_checks_misc
from wordless.wl_dialogs import (
    wl_dialogs,
    wl_dialogs_errs,
    wl_dialogs_misc,
    wl_msg_boxes
)
from wordless.wl_nlp import wl_matching, wl_nlp_utils, wl_texts
from wordless.wl_utils import (
    wl_conversion,
    wl_detection,
    wl_misc,
    wl_paths,
    wl_threading
)
from wordless.wl_widgets import (
    wl_boxes,
    wl_buttons,
    wl_item_delegates,
    wl_labels,
    wl_layouts,
    wl_tables
)

_tr = QCoreApplication.translate

class Wrapper_File_Area(wl_layouts.Wl_Wrapper_File_Area):
    def __init__(self, main, file_type = 'observed'):
        super().__init__(main)

        self.file_names_old = []
        self.file_type = file_type

        # Suffix for settings
        if self.file_type == 'observed':
            self.tab = 'corpora_observed'
            self.settings_suffix = ''
        elif self.file_type == 'ref':
            self.tab = 'corpora_ref'
            self.settings_suffix = '_ref'

        # Table
        self.table_files = Wl_Table_Files(self)

        self.wrapper_table.layout().addWidget(self.table_files, 0, 0)

        # Load files
        self.table_files.update_table()

    def get_files(self):
        return self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}']

    def get_file_names(self):
        return (
            file['name']
            for file in self.get_files()
        )

    def get_selected_files(self):
        return (
            file
            for file in self.get_files()
            if file['selected']
        )

    def get_selected_file_names(self):
        return (
            file['name']
            for file in self.get_selected_files()
        )

    def find_file_by_name(self, file_name, selected_only = False):
        if selected_only:
            files = self.get_selected_files()
        else:
            files = self.get_files()

        for file in files:
            if file['name'] == file_name:
                return file

        return None

    def find_files_by_name(self, file_names, selected_only = False):
        files = [
            self.find_file_by_name(file_name, selected_only = selected_only)
            for file_name in file_names
        ]

        return (file for file in files if file)

# References:
#     https://stackoverflow.com/a/29621256
#     https://wiki.qt.io/Technical_FAQ#How_can_I_insert_a_checkbox_into_the_header_of_my_view.3F
class Wl_Table_Header_Files(QHeaderView):
    def __init__(self, orientation, parent):
        super().__init__(orientation, parent)

        self.table = parent
        self._is_checked = False

        self.setSectionsClickable(True)

        self.sectionClicked.connect(self.section_clicked)

    def paintSection(self, painter, rect, logicalIndex):
        painter.save()

        super().paintSection(painter, rect, logicalIndex)

        painter.restore()

        if logicalIndex == 0:
            option = QStyleOptionButton()
            option.rect = QRect(3, 6, 16, 16)

            if self.table.is_empty():
                option.state = QStyle.State_None
            else:
                option.state = QStyle.State_Enabled | QStyle.State_Active

            if self._is_checked:
                option.state |= QStyle.State_On
            else:
                option.state |= QStyle.State_Off

            self.style().drawPrimitive(QStyle.PE_IndicatorCheckBox, option, painter)

    def section_clicked(self, logicalIndex):
        if logicalIndex == 0:
            self._is_checked = not self._is_checked

            if self._is_checked:
                self.select_all()
            else:
                self.deselect_all()

            self.update()

    def select_all(self):
        self._is_checked = True

        for i in range(self.model().rowCount()):
            self.model().item(i, 0).setCheckState(Qt.Checked)

    def deselect_all(self):
        self._is_checked = True

        for i in range(self.model().rowCount()):
            self.model().item(i, 0).setCheckState(Qt.Unchecked)

    def invert_selection(self):
        for i in range(self.model().rowCount()):
            if self.model().item(i, 0).checkState() == Qt.Checked:
                self.model().item(i, 0).setCheckState(Qt.Unchecked)
            else:
                self.model().item(i, 0).setCheckState(Qt.Checked)

class Wl_Table_Files(wl_tables.Wl_Table):
    def __init__(self, parent):
        super().__init__(
            parent,
            headers = [
                # Padding for the checkbox
                _tr('Wl_Table_Files', '  Name'),
                _tr('Wl_Table_Files', 'Path'),
                _tr('Wl_Table_Files', 'Encoding'),
                _tr('Wl_Table_Files', 'Language'),
                _tr('Wl_Table_Files', 'Tokenized'),
                _tr('Wl_Table_Files', 'Tagged')
            ],
            editable = True,
            drag_drop = True
        )

        self.file_area = parent
        self.file_type = self.file_area.file_type
        self.settings_suffix = self.file_area.settings_suffix

        self.setHorizontalHeader(Wl_Table_Header_Files(Qt.Horizontal, self))

        self.setItemDelegateForColumn(1, wl_item_delegates.Wl_Item_Delegate_Uneditable(self))
        self.setItemDelegateForColumn(2, wl_item_delegates.Wl_Item_Delegate_Uneditable(self))
        self.setItemDelegateForColumn(3, wl_item_delegates.Wl_Item_Delegate_Uneditable(self))
        self.setItemDelegateForColumn(4, wl_item_delegates.Wl_Item_Delegate_Uneditable(self))
        self.setItemDelegateForColumn(5, wl_item_delegates.Wl_Item_Delegate_Uneditable(self))

        self.selectionModel().selectionChanged.connect(self.selection_changed)
        self.clicked.connect(self.item_clicked)

        # Menu
        self.main.action_file_open_files.triggered.connect(
            lambda: self.check_file_area(self.open_files)
        )
        self.main.action_file_reopen.triggered.connect(
            lambda: self.check_file_area(self.reopen)
        )

        self.main.action_file_select_all.triggered.connect(
            lambda: self.check_file_area(self.horizontalHeader().select_all)
        )
        self.main.action_file_deselect_all.triggered.connect(
            lambda: self.check_file_area(self.horizontalHeader().deselect_all)
        )
        self.main.action_file_invert_selection.triggered.connect(
            lambda: self.check_file_area(self.horizontalHeader().invert_selection)
        )

        self.main.action_file_close_selected.triggered.connect(
            lambda: self.check_file_area(self.close_selected)
        )
        self.main.action_file_close_all.triggered.connect(
            lambda: self.check_file_area(self.close_all)
        )

        self.main.tabs_file_area.currentChanged.connect(
            lambda: self.check_file_area(self.model().itemChanged.emit, self.model().item(0, 0))
        )

    def item_changed(self):
        super().item_changed()

        if not self.is_empty():
            # Record old file names that might be useful for other slots
            self.file_area.file_names_old = list(self.file_area.get_selected_file_names())

            # Check for empty and duplicate file names
            for row in range(self.model().rowCount()):
                file = self.model().item(row, 0).wl_file
                file_name = self.model().item(row, 0).text()

                if file_name != file['name_old']:
                    if not file_name or self.main.wl_file_area.find_file_by_name(file_name):
                        self.disable_updates()

                        self.model().item(row, 0).setText(file['name_old'])

                        self.enable_updates()

                        if not file_name:
                            wl_msg_boxes.Wl_Msg_Box_Warning(
                                self.main,
                                title = self.tr('Empty File Name'),
                                text = self.tr('''
                                    <div>The file name should not be left empty!</div>
                                ''')
                            ).exec_()
                        elif self.main.wl_file_area.find_file_by_name(file_name):
                            wl_msg_boxes.Wl_Msg_Box_Warning(
                                self.main,
                                title = self.tr('Duplicate File Names'),
                                text = self.tr('''
                                    <div>There is already a file with the same name in the file area.</div>
                                    <div>Please specify a different file name.</div>
                                ''')
                            ).exec_()

                        self.setCurrentIndex(self.model().index(row, 0))

                        self.closeEditor(self.findChild(QLineEdit), QAbstractItemDelegate.NoHint)
                        self.edit(self.model().index(row, 0))

                    break

            self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}'].clear()

            for row in range(self.model().rowCount()):
                file = self.model().item(row, 0).wl_file

                file['selected'] = self.model().item(row, 0).checkState() == Qt.Checked
                file['name'] = file['name_old'] = self.model().item(row, 0).text()
                file['encoding'] = wl_conversion.to_encoding_code(self.main, self.model().item(row, 2).text())
                file['lang'] = wl_conversion.to_lang_code(self.main, self.model().item(row, 3).text())
                file['tokenized'] = wl_conversion.to_yes_no_code(self.model().item(row, 4).text())
                file['tagged'] = wl_conversion.to_yes_no_code(self.model().item(row, 5).text())

                self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}'].append(file)

            # Checkboxes
            check_states = []

            for i in range(self.model().rowCount()):
                if self.model().item(i, 0).checkState() == Qt.Checked:
                    check_states.append(Qt.Checked)
                else:
                    check_states.append(Qt.Unchecked)

            if all((check_state == Qt.Checked for check_state in check_states)):
                self.horizontalHeader()._is_checked = True
            else:
                self.horizontalHeader()._is_checked = False

        self.horizontalHeader().update()

        # Menu
        if not self.is_empty():
            self.main.action_file_select_all.setEnabled(True)
            self.main.action_file_deselect_all.setEnabled(True)
            self.main.action_file_invert_selection.setEnabled(True)

            self.main.action_file_close_all.setEnabled(True)
        else:
            self.main.action_file_select_all.setEnabled(False)
            self.main.action_file_deselect_all.setEnabled(False)
            self.main.action_file_invert_selection.setEnabled(False)

            self.main.action_file_close_all.setEnabled(False)

        if self.main.settings_custom['file_area'][f'files_closed{self.settings_suffix}']:
            self.main.action_file_reopen.setEnabled(True)
        else:
            self.main.action_file_reopen.setEnabled(False)

        self.selectionModel().selectionChanged.emit(QItemSelection(), QItemSelection())

    def item_clicked(self):
        if not self.is_empty():
            for row in range(self.model().rowCount()):
                if self.model().item(row, 0).checkState() == Qt.Checked:
                    self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}'][row]['selected'] = True
                else:
                    self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}'][row]['selected'] = False

    def selection_changed(self):
        if self.get_selected_rows():
            self.main.action_file_close_selected.setEnabled(True)
        else:
            self.main.action_file_close_selected.setEnabled(False)

    def update_table(self):
        if (files := self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}']):
            self.clr_table(len(files))

            self.disable_updates()

            for i, file in enumerate(files):
                item_name = QStandardItem(file['name'])
                # Record file properties
                item_name.wl_file = file
                item_name.setCheckable(True)

                if file['selected']:
                    item_name.setCheckState(Qt.Checked)
                else:
                    item_name.setCheckState(Qt.Unchecked)

                self.model().setItem(i, 0, item_name)
                self.model().setItem(i, 1, QStandardItem(file['path_orig']))
                self.model().setItem(i, 2, QStandardItem(wl_conversion.to_encoding_text(self.main, file['encoding'])))
                self.model().setItem(i, 3, QStandardItem(wl_conversion.to_lang_text(self.main, file['lang'])))
                self.model().setItem(i, 4, QStandardItem(wl_conversion.to_yes_no_text(file['tokenized'])))
                self.model().setItem(i, 5, QStandardItem(wl_conversion.to_yes_no_text(file['tagged'])))

            self.enable_updates()
        else:
            self.clr_table(1)

    def check_file_area(self, op, *args, **kwargs):
        if (
            (
                self.file_type == 'observed'
                and self.main.tabs_file_area.tabText(self.main.tabs_file_area.currentIndex()) == self.tr('Observed Corpora')
            ) or (
                self.file_type == 'ref'
                and self.main.tabs_file_area.tabText(self.main.tabs_file_area.currentIndex()) == self.tr('Reference Corpora')
            )
        ):
            return op(*args, **kwargs)

        return None

    @wl_misc.log_time
    def _open_files(self, files_to_open):
        if wl_nlp_utils.check_models(
            self.main,
            langs = set((file['lang'] for file in files_to_open)),
        ):
            dialog_progress = wl_dialogs_misc.Wl_Dialog_Progress(self.main, text = self.tr('Checking files...'))

            wl_threading.Wl_Thread(Wl_Worker_Open_Files(
                self.main,
                dialog_progress = dialog_progress,
                update_gui = self.update_gui,
                files_to_open = files_to_open,
                file_type = self.file_type
            )).start_worker()

    def update_gui(self, err_msg, new_files):
        if wl_checks_files.check_err_file_area(self.main, err_msg):
            len_files_old = len(self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}'])

            self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}'].extend(new_files)
            self.update_table()

            len_files_opened = len(self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}']) - len_files_old
            msg_file = self.tr('file') if len_files_opened == 1 else self.tr('files')

            self.main.statusBar().showMessage(self.tr('{} {} has been successfully opened.').format(len_files_opened, msg_file))

    def open_files(self):
        self.dialog_open_files = Wl_Dialog_Open_Files(self.main)
        self.dialog_open_files.open()

    def reopen(self):
        files = self.main.settings_custom['file_area'][f'files_closed{self.settings_suffix}'].pop()

        dialog_open_files = Wl_Dialog_Open_Files(self.main)
        dialog_open_files._add_files(list(dict.fromkeys([file['path_orig'] for file in files])))

        self._open_files(files_to_open = dialog_open_files.table_files.files_to_open)

    def _close_files(self, i_files):
        self.main.settings_custom['file_area'][f'files_closed{self.settings_suffix}'].append([])

        for i in reversed(i_files):
            file_to_remove = self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}'].pop(i)

            self.main.settings_custom['file_area'][f'files_closed{self.settings_suffix}'][-1].append(file_to_remove)

            # Remove temporary files
            if os.path.exists(file_to_remove['path']):
                os.remove(file_to_remove['path'])

        self.update_table()

    def close_selected(self):
        self._close_files(self.get_selected_rows())

    def close_all(self):
        self._close_files(list(range(len(self.main.settings_custom['file_area'][f'files_open{self.settings_suffix}']))))

class Wl_Dialog_Open_Files(wl_dialogs.Wl_Dialog):
    def __init__(self, main):
        super().__init__(
            main,
            title = _tr('Wl_Dialog_Open_Files', 'Open Files'),
            width = 800,
            height = 320
        )

        self.table_files = Table_Open_Files(self)

        self.table_files.model().itemChanged.connect(self.table_files_changed)

        self.table_files.button_add.hide()
        self.table_files.button_ins.hide()

        self.table_files.button_add_files = QPushButton(self.tr('Add files...'), self)
        self.table_files.button_add_folder = QPushButton(self.tr('Add folder...'), self)
        self.table_files.button_del.setText(self.tr('Remove files'))
        self.table_files.button_clr.setText(self.tr('Clear table'))

        self.table_files.button_add_files.setMinimumWidth(120)
        self.table_files.button_add_files.setMinimumWidth(120)
        self.table_files.button_add_files.setMinimumWidth(120)
        self.table_files.button_add_files.setMinimumWidth(120)

        self.table_files.button_add_files.clicked.connect(self.add_files)
        self.table_files.button_add_folder.clicked.connect(self.add_folder)

        layout_table = wl_layouts.Wl_Layout()
        layout_table.addWidget(self.table_files, 0, 0, 5, 1)
        layout_table.addWidget(self.table_files.button_add_files, 0, 1)
        layout_table.addWidget(self.table_files.button_add_folder, 1, 1)
        layout_table.addWidget(self.table_files.button_del, 2, 1)
        layout_table.addWidget(self.table_files.button_clr, 3, 1)

        layout_table.setRowStretch(4, 1)

        self.checkbox_auto_detect_encodings = QCheckBox(self.tr('Auto-detect encodings'), self)
        self.checkbox_auto_detect_langs = QCheckBox(self.tr('Auto-detect languages'), self)
        self.checkbox_include_files_in_subfolders = QCheckBox(self.tr('Include files in subfolders'), self)

        self.checkbox_auto_detect_encodings.stateChanged.connect(self.settings_changed)
        self.checkbox_auto_detect_langs.stateChanged.connect(self.settings_changed)
        self.checkbox_include_files_in_subfolders.stateChanged.connect(self.settings_changed)

        layout_checkboxes = wl_layouts.Wl_Layout()
        layout_checkboxes.addWidget(self.checkbox_auto_detect_encodings, 0, 0)
        layout_checkboxes.addWidget(self.checkbox_auto_detect_langs, 0, 1)
        layout_checkboxes.addWidget(self.checkbox_include_files_in_subfolders, 1, 0)

        self.button_restore_defaults = wl_buttons.Wl_Button_Restore_Defaults(self, load_settings = self.load_settings)
        self.button_open = QPushButton(self.tr('Open'), self)
        self.button_cancel = QPushButton(self.tr('Cancel'), self)

        self.button_open.clicked.connect(self.accept)
        self.button_cancel.clicked.connect(self.reject)

        self.setLayout(wl_layouts.Wl_Layout())
        self.layout().addLayout(layout_table, 0, 0, 1, 4)
        self.layout().addLayout(layout_checkboxes, 1, 0, 1, 4)

        self.layout().addWidget(wl_layouts.Wl_Separator(self), 2, 0, 1, 4)

        self.layout().addWidget(self.button_restore_defaults, 3, 0)
        self.layout().addWidget(self.button_open, 3, 2)
        self.layout().addWidget(self.button_cancel, 3, 3)

        self.layout().setColumnStretch(1, 1)

        self.load_settings()

    def accept(self):
        num_files = len(self.main.settings_custom['file_area']['files_open'] + self.main.settings_custom['file_area']['files_open_ref'])

        self.main.tabs_file_area.currentWidget().table_files._open_files(files_to_open = self.table_files.files_to_open)

        if num_files < len(self.main.settings_custom['file_area']['files_open'] + self.main.settings_custom['file_area']['files_open_ref']):
            super().accept()

    def reject(self):
        # Remove placeholders for new paths
        for file in self.table_files.files_to_open:
            if os.path.exists(file['path']):
                os.remove(file['path'])

        super().reject()

    def load_settings(self, defaults = False):
        if defaults:
            settings = copy.deepcopy(self.main.settings_default['file_area']['dialog_open_files'])
        else:
            settings = copy.deepcopy(self.main.settings_custom['file_area']['dialog_open_files'])

        self.checkbox_auto_detect_encodings.setChecked(settings['auto_detect_encodings'])
        self.checkbox_auto_detect_langs.setChecked(settings['auto_detect_langs'])
        self.checkbox_include_files_in_subfolders.setChecked(settings['include_files_in_subfolders'])

        self.table_files.model().itemChanged.emit(QStandardItem())
        self.settings_changed()

    def table_files_changed(self, item): # pylint: disable=unused-argument
        if self.table_files.is_empty():
            self.button_open.setEnabled(False)
        else:
            self.button_open.setEnabled(True)

    def settings_changed(self):
        settings = self.main.settings_custom['file_area']['dialog_open_files']

        settings['auto_detect_encodings'] = self.checkbox_auto_detect_encodings.isChecked()
        settings['auto_detect_langs'] = self.checkbox_auto_detect_langs.isChecked()
        settings['include_files_in_subfolders'] = self.checkbox_include_files_in_subfolders.isChecked()

    def _add_files(self, file_paths):
        dialog_progress = wl_dialogs_misc.Wl_Dialog_Progress(self.main, text = self.tr('Checking files...'))

        file_paths, self.file_paths_unsupported = wl_checks_files.check_file_paths_unsupported(self.main, file_paths)
        file_paths, self.file_paths_empty = wl_checks_files.check_file_paths_empty(self.main, file_paths)
        file_paths, self.file_paths_dup = wl_checks_files.check_file_paths_dup(
            self.main,
            new_file_paths = file_paths,
            file_paths = [
                file['path_orig']
                for file in (
                    self.main.settings_custom['file_area'][f'files_open{self.main.tabs_file_area.currentWidget().settings_suffix}']
                    + self.table_files.files_to_open
                )
            ]
        )

        # Display warning when opening non-text files
        if (
            any((
                os.path.splitext(file_path)[1].lower() not in ['.csv', '.lrc', '.txt', '.tmx', '.xml']
                for file_path in file_paths
            ))
            and self.main.settings_custom['files']['misc_settings']['display_warning_when_opening_nontext_files']
        ):
            non_text_files_ok = Wl_Dialog_Opening_Nontext_Files(self.main).exec_()
        else:
            non_text_files_ok = True

        if non_text_files_ok:
            wl_threading.Wl_Thread(Wl_Worker_Add_Files(
                self.main,
                dialog_progress = dialog_progress,
                update_gui = self.update_gui,
                file_paths = file_paths,
                table = self.table_files
            )).start_worker()

    def update_gui(self, err_msg, new_files):
        if wl_checks_files.check_err_file_area(self.main, err_msg):
            self.table_files.files_to_open.extend(new_files)

            self.table_files.update_table()

            if self.file_paths_empty or self.file_paths_unsupported or self.file_paths_dup:
                dialog_err_files = wl_dialogs_errs.Wl_Dialog_Err_Files(self.main, title = self.tr('Error Adding Files'))

                dialog_err_files.label_err.set_text(self.tr('''
                    <div>
                        An error occurred while adding files, so the following files are not added to the table.
                    </div>
                '''))
                dialog_err_files.table_err_files.model().setRowCount(
                    len(self.file_paths_empty)
                    + len(self.file_paths_unsupported)
                    + len(self.file_paths_dup)
                )

                dialog_err_files.table_err_files.disable_updates()

                for i, file_path in enumerate(self.file_paths_empty + self.file_paths_unsupported + self.file_paths_dup):
                    if file_path in self.file_paths_empty:
                        dialog_err_files.table_err_files.model().setItem(
                            i, 0,
                            QStandardItem(self.tr('Empty file'))
                        )
                    elif file_path in self.file_paths_unsupported:
                        dialog_err_files.table_err_files.model().setItem(
                            i, 0,
                            QStandardItem(self.tr('Unsupported file type'))
                        )
                    elif file_path in self.file_paths_dup:
                        dialog_err_files.table_err_files.model().setItem(
                            i, 0,
                            QStandardItem(self.tr('Duplicate file'))
                        )

                    dialog_err_files.table_err_files.model().setItem(
                        i, 1,
                        QStandardItem(file_path)
                    )

                dialog_err_files.table_err_files.enable_updates()
                dialog_err_files.exec_()

    def add_files(self):
        if os.path.exists(self.main.settings_custom['general']['imp']['files']['default_path']):
            default_dir = self.main.settings_custom['general']['imp']['files']['default_path']
        else:
            default_dir = self.main.settings_default['general']['imp']['files']['default_path']

        file_paths = QFileDialog.getOpenFileNames(
            parent = self.main,
            caption = self.tr('Open Files'),
            directory = wl_checks_misc.check_dir(default_dir),
            filter = ';;'.join(self.main.settings_global['file_types']['files']),
            initialFilter = self.main.settings_global['file_types']['files'][-1]
        )[0]

        if file_paths:
            self._add_files(file_paths)

    def add_folder(self):
        file_paths = []

        file_dir = QFileDialog.getExistingDirectory(
            parent = self.main,
            caption = self.tr('Open Folder'),
            directory = self.main.settings_custom['general']['imp']['files']['default_path']
        )

        if file_dir:
            if self.main.settings_custom['file_area']['dialog_open_files']['include_files_in_subfolders']:
                for dir_path, _, file_names in os.walk(file_dir):
                    for file_name in file_names:
                        file_paths.append(os.path.join(dir_path, file_name))
            else:
                file_names = list(os.walk(file_dir))[0][2]

                for file_name in file_names:
                    file_paths.append(os.path.join(file_dir, file_name))

            self._add_files(file_paths)

class Table_Open_Files(wl_tables.Wl_Table_Add_Ins_Del_Clr):
    def __init__(self, parent):
        super().__init__(
            parent = parent,
            headers = [
                _tr('Table_Open_Files', 'Path'),
                _tr('Table_Open_Files', 'Encoding'),
                _tr('Table_Open_Files', 'Language'),
                _tr('Table_Open_Files', 'Tokenized'),
                _tr('Table_Open_Files', 'Tagged')
            ],
            col_edit = 2
        )

        self.files_to_open = []

        self.setItemDelegateForColumn(0, wl_item_delegates.Wl_Item_Delegate_Uneditable(self))
        self.setItemDelegateForColumn(1, wl_item_delegates.Wl_Item_Delegate_Combo_Box_Custom(self, wl_boxes.Wl_Combo_Box_Encoding))
        self.setItemDelegateForColumn(2, wl_item_delegates.Wl_Item_Delegate_Combo_Box_Custom(self, wl_boxes.Wl_Combo_Box_Lang))
        self.setItemDelegateForColumn(3, wl_item_delegates.Wl_Item_Delegate_Combo_Box_Custom(self, wl_boxes.Wl_Combo_Box_Yes_No))
        self.setItemDelegateForColumn(4, wl_item_delegates.Wl_Item_Delegate_Combo_Box_Custom(self, wl_boxes.Wl_Combo_Box_Yes_No))

        self.button_clr.disconnect()
        self.button_clr.clicked.connect(lambda: self.clr_table(remove_placeholders = True))

        self.clr_table()

    def item_changed(self):
        super().item_changed()

        self.files_to_open = []

        if not self.is_empty():
            for row in range(self.model().rowCount()):
                file = self.model().item(row, 0).file

                file['encoding'] = wl_conversion.to_encoding_code(self.main, self.model().item(row, 1).text())
                file['lang'] = wl_conversion.to_lang_code(self.main, self.model().item(row, 2).text())
                file['tokenized'] = wl_conversion.to_yes_no_code(self.model().item(row, 3).text())
                file['tagged'] = wl_conversion.to_yes_no_code(self.model().item(row, 4).text())

                self.files_to_open.append(file)

    def del_row(self):
        for row in self.get_selected_rows():
            file_path = self.files_to_open[row]['path']

            if os.path.exists(file_path):
                os.remove(file_path)

        super().del_row()

    def clr_table(self, num_headers = 1, remove_placeholders = False):
        # Remove placeholders for new paths
        if remove_placeholders:
            for file in self.files_to_open:
                if os.path.exists(file['path']):
                    os.remove(file['path'])

        super().clr_table(num_headers = num_headers)

    def update_table(self):
        files = self.files_to_open

        if files:
            self.clr_table(len(files))

            self.disable_updates()

            for i, file in enumerate(files):
                self.model().setItem(i, 0, QStandardItem(file['path_orig']))
                self.model().setItem(i, 1, QStandardItem(wl_conversion.to_encoding_text(self.main, file['encoding'])))
                self.model().setItem(i, 2, QStandardItem(wl_conversion.to_lang_text(self.main, file['lang'])))
                self.model().setItem(i, 3, QStandardItem(wl_conversion.to_yes_no_text(file['tokenized'])))
                self.model().setItem(i, 4, QStandardItem(wl_conversion.to_yes_no_text(file['tagged'])))

                self.model().item(i, 0).file = file

            self.enable_updates()
        else:
            self.clr_table()

class Wl_Dialog_Opening_Nontext_Files(wl_dialogs.Wl_Dialog_Info):
    def __init__(self, main):
        super().__init__(
            main,
            title = _tr('Wl_Dialog_Opening_Nontext_Files', 'Opening Non-text Files'),
            width = 550,
            no_buttons = True
        )

        self.label_opening_non_text_files = wl_labels.Wl_Label_Dialog(
            self.tr('''
                <div>It is <b>not recommended to directly import non-text files into <i>Wordless</i></b> and the support for doing so is provided only for convenience, since accuracy of text extraction could never be guaranteed and unintended data loss might occur, for which reason users are encouraged to <b>convert their files using specialized tools and make their own choices</b> on which part of the data should be kept or discarded.</div>
                <br>
                <div>Do you want to proceed to open non-text files anyway?</div>
            '''),
            self
        )

        self.checkbox_do_not_show_this_again = QCheckBox(self.tr('Do not show this again'), self)
        self.button_proceed = QPushButton(self.tr('Proceed'), self)
        self.button_abort = QPushButton(self.tr('Abort'), self)

        self.checkbox_do_not_show_this_again.stateChanged.connect(self.do_not_show_this_again_changed)
        self.button_proceed.clicked.connect(self.accept)
        self.button_abort.clicked.connect(self.reject)

        self.layout_info.addWidget(self.label_opening_non_text_files, 0, 0)

        self.layout_buttons.addWidget(self.checkbox_do_not_show_this_again, 0, 0)
        self.layout_buttons.addWidget(self.button_proceed, 0, 2)
        self.layout_buttons.addWidget(self.button_abort, 0, 3)

        self.layout_buttons.setColumnStretch(1, 1)

        self.load_settings()

    def load_settings(self):
        settings = copy.deepcopy(self.main.settings_custom['files']['misc_settings'])

        self.checkbox_do_not_show_this_again.setChecked(not settings['display_warning_when_opening_nontext_files'])

    def do_not_show_this_again_changed(self):
        settings = self.main.settings_custom['files']['misc_settings']

        settings['display_warning_when_opening_nontext_files'] = not self.checkbox_do_not_show_this_again.isChecked()

# Reference: https://github.com/python-openxml/python-docx/issues/40#issuecomment-1793226714
def iter_block_items(blkcntnr):
    for item in blkcntnr.iter_inner_content():
        if isinstance(item, docx.text.paragraph.Paragraph):
            yield item
        elif isinstance(item, docx.table.Table):
            for row in iter_visual_cells(item):
                for cell in row:
                    yield from iter_block_items(cell)

# Reference: https://github.com/python-openxml/python-docx/issues/344#issuecomment-271390490
def iter_visual_cells(table):
    visual_cells = []
    prior_tcs = set()

    for row in table.rows:
        visual_cells.append([])

        for cell in row.cells:
            if cell._tc in prior_tcs: # skip cells pointing to same `<w:tc>` element
                continue
            else:
                visual_cells[-1].append(cell)

                prior_tcs.add(cell._tc)

    return visual_cells

# Reference: https://stackoverflow.com/questions/51701626/how-to-extract-text-from-a-text-shape-within-a-group-shape-in-powerpoint-using
def iter_slide_shapes(shapes):
    texts = []

    for shape in shapes:
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP: # pylint: disable=no-member
            iter_slide_shapes(shape)

        if shape.has_text_frame:
            texts.append(shape.text)

    return texts

LRC_TIME_TAGS_VALID = r'[0-9]{2}:[0-5][0-9][\.:][0-9]{2,3}'
RE_LRC_TIME_TAGS_LINE_START = re.compile(r'^\[[^\]]+?\]')
RE_LRC_TIME_TAGS_VALID = re.compile(fr'^\[{LRC_TIME_TAGS_VALID}\]$')
RE_LRC_TIME_TAGS_WORDS = re.compile(fr'\<{LRC_TIME_TAGS_VALID}\>')

def get_text_non_tmx(file):
    file_path = file['path_orig']
    file_ext = os.path.splitext(os.path.basename(file_path))[1].lower()

    match file_ext:
        # Text and XML files
        case '.txt' | '.xml':
            with open(file_path, 'r', encoding = file['encoding'], errors = 'replace') as f:
                text = f.read()
        # CSV files
        case '.csv':
            lines = []

            with open(file_path, 'r', encoding = file['encoding'], errors = 'replace', newline = '') as f:
                # Remove NULL bytes to avoid error
                csv_reader = csv.reader([line.replace('\0', '') for line in f])

                for row in csv_reader:
                    lines.append('\t'.join(row))

            text = '\n'.join(lines)
        # Excel workbooks
        case '.xlsx':
            lines = []
            workbook = openpyxl.load_workbook(file_path, data_only = True)

            for worksheet_name in workbook.sheetnames:
                worksheet = workbook[worksheet_name]

                for row in worksheet.rows:
                    cells = [
                        # Numbers need to be converted to strings
                        (str(cell.value) if cell.value is not None else '')
                        for cell in row
                    ]

                    lines.append('\t'.join(cells))

            text = '\n'.join(lines)
        # HTML pages
        case '.htm' | '.html':
            with open(file_path, 'r', encoding = file['encoding'], errors = 'replace') as f:
                soup = bs4.BeautifulSoup(f.read(), 'lxml')

            text = soup.get_text()
        # Lyrics files
        case '.lrc':
            lyrics = {}

            with open(file_path, 'r', encoding = file['encoding'], errors = 'replace') as f:
                for line in f:
                    time_tags = []

                    line = line.strip()

                    # Extract time tags at the beginning of the line
                    while (re_time_tag := RE_LRC_TIME_TAGS_LINE_START.search(line)):
                        time_tags.append(re_time_tag.group())

                        line = line[len(time_tags[-1]):].strip()

                    # Strip word time tags
                    line = RE_LRC_TIME_TAGS_WORDS.sub(r'', line)
                    line = re.sub(r'\s{2,}', r' ', line).strip()

                    for time_tag in time_tags:
                        if RE_LRC_TIME_TAGS_VALID.search(time_tag):
                            lyrics[time_tag] = line

            text = '\n'.join((lyrics_line for _, lyrics_line in sorted(lyrics.items()))) + '\n'
        # PDF files
        case '.pdf':
            reader = pypdf.PdfReader(file_path)
            text = '\n'.join([page.extract_text() for page in reader.pages])
        # PowerPoint presentations
        case '.pptx':
            texts = []
            prs = pptx.Presentation(file_path)

            for slide in prs.slides:
                texts.extend(iter_slide_shapes(slide.shapes))

            text = '\n'.join(texts)
        # Word documents
        # Reference: https://github.com/python-openxml/python-docx/issues/40#issuecomment-1793226714
        case '.docx':
            lines = []
            doc = docx.Document(file_path)

            for item in doc.iter_inner_content():
                if isinstance(item, docx.text.paragraph.Paragraph):
                    lines.append(item.text)
                elif isinstance(item, docx.table.Table):
                    for row in iter_visual_cells(item):
                        cells = [
                            ' '.join([cell_item.text for cell_item in iter_block_items(cell)])
                            for cell in row
                        ]

                        lines.append('\t'.join(cells))

            text = '\n'.join(lines)

    return text

class Wl_Worker_Add_Files(wl_threading.Wl_Worker):
    worker_done = pyqtSignal(str, list)

    def run(self):
        err_msg = ''
        new_files = []

        try:
            len_file_paths = len(self.file_paths)

            for i, file_path in enumerate(self.file_paths):
                self.progress_updated.emit(self.tr('Adding files... ({}/{})').format(i + 1, len_file_paths))

                file_path = wl_paths.get_normalized_path(file_path)
                file_name, file_ext = os.path.splitext(os.path.basename(file_path))
                file_ext = file_ext.lower()

                new_file = {'selected': True, 'path_orig': file_path}

                # Check for duplicate file names
                file_names = [
                    *self.main.wl_file_area.get_file_names(),
                    *[file['name'] for file in self.table.files_to_open],
                    *[new_file['name'] for new_file in new_files]
                ]

                new_file['name'] = new_file['name_old'] = wl_checks_misc.check_new_name(file_name, file_names)

                # Path, Tokenized, Tagged
                default_dir = wl_checks_misc.check_dir(self.main.settings_custom['general']['imp']['temp_files']['default_path'])

                if file_ext == '.xml':
                    new_file['path'] = os.path.join(default_dir, f'{file_name}.xml')

                    # Use default settings for "Tokenized" & "Tagged" if auto-detection of encodings and languages are both disabled
                    if (
                        not self.main.settings_custom['file_area']['dialog_open_files']['auto_detect_encodings']
                        and not self.main.settings_custom['file_area']['dialog_open_files']['auto_detect_langs']
                    ):
                        new_file['tokenized'] = self.main.settings_custom['files']['default_settings']['tokenized']
                        new_file['tagged'] = self.main.settings_custom['files']['default_settings']['tagged']
                    else:
                        new_file['tokenized'] = True
                        new_file['tagged'] = True
                else:
                    new_file['path'] = os.path.join(default_dir, f'{file_name}.txt')
                    new_file['tokenized'] = self.main.settings_custom['files']['default_settings']['tokenized']
                    new_file['tagged'] = self.main.settings_custom['files']['default_settings']['tagged']

                # Check for duplicate files
                new_file['path'] = wl_checks_misc.check_new_path(new_file['path'])

                # Detect encodings
                default_encoding = self.main.settings_custom['files']['default_settings']['encoding']

                if file_ext in ['.docx', '.xlsx']:
                    new_file['encoding'] = default_encoding
                else:
                    if self.main.settings_custom['file_area']['dialog_open_files']['auto_detect_encodings']:
                        new_file['encoding'] = wl_detection.detect_encoding(self.main, file_path)
                    else:
                        new_file['encoding'] = default_encoding

                # Cleanse contents before language detection
                if file_ext != '.tmx':
                    new_file['text'] = get_text_non_tmx(new_file)

                    if self.main.settings_custom['file_area']['dialog_open_files']['auto_detect_langs']:
                        new_file['lang'] = wl_detection.detect_lang_text(self.main, new_file['text'])
                    else:
                        new_file['lang'] = self.main.settings_custom['files']['default_settings']['lang']

                    new_files.append(new_file)
                # Translation memory files
                else:
                    lines_src = []
                    lines_tgt = []

                    new_file_src = copy.deepcopy(new_file)
                    new_file_tgt = copy.deepcopy(new_file)

                    new_file_src['tmx_type'] = 'src'
                    new_file_tgt['tmx_type'] = 'tgt'

                    new_file_src['name'] = new_file_src['name_old'] = wl_checks_misc.check_new_name(f'{file_name}_source', file_names)
                    new_file_tgt['name'] = new_file_tgt['name_old'] = wl_checks_misc.check_new_name(f'{file_name}_target', file_names)

                    new_file_src['path'] = wl_checks_misc.check_new_path(os.path.join(default_dir, f'{file_name}_source.txt'))
                    new_file_tgt['path'] = wl_checks_misc.check_new_path(os.path.join(default_dir, f'{file_name}_target.txt'))

                    with open(file_path, 'r', encoding = new_file['encoding'], errors = 'replace') as f:
                        soup = bs4.BeautifulSoup(f.read(), 'lxml-xml')

                    # Identify source and target languages
                    elements_tuv = soup.select(r'tu:first-child tuv[xml\:lang]')

                    if len(elements_tuv) == 2:
                        new_file_src['lang'] = wl_conversion.to_iso_639_3(self.main, elements_tuv[0]['xml:lang'])
                        new_file_tgt['lang'] = wl_conversion.to_iso_639_3(self.main, elements_tuv[1]['xml:lang'])

                        if new_file_src['lang'] is None:
                            new_file_src['lang'] = 'other'
                        if new_file_tgt['lang'] is None:
                            new_file_tgt['lang'] = 'other'
                    else:
                        new_file_src['lang'] = new_file_tgt['lang'] = self.main.settings_custom['files']['default_settings']['lang']

                    with open(new_file['path_orig'], 'r', encoding = new_file['encoding'], errors = 'replace') as f:
                        soup = bs4.BeautifulSoup(f.read(), 'lxml-xml')

                    for elements_tu in soup.select('tu'):
                        seg_src, seg_tgt = elements_tu.select('seg')

                        lines_src.append(seg_src.get_text().replace(r'\n', ' ').strip())
                        lines_tgt.append(seg_tgt.get_text().replace(r'\n', ' ').strip())

                    new_file_src['text'] = '\n'.join(lines_src)
                    new_file_tgt['text'] = '\n'.join(lines_tgt)

                    new_files.append(new_file_src)
                    new_files.append(new_file_tgt)

            if self.file_paths:
                self.main.settings_custom['general']['imp']['files']['default_path'] = wl_paths.get_normalized_dir(self.file_paths[0])
        except Exception:
            err_msg = traceback.format_exc()

        self.progress_updated.emit(self.tr('Updating table...'))
        self.worker_done.emit(err_msg, new_files)

class Wl_Worker_Open_Files(wl_threading.Wl_Worker):
    worker_done = pyqtSignal(str, list)

    def run(self):
        err_msg = ''
        new_files = []

        try:
            len_files = len(self.files_to_open)
            # Regex for headers
            tags_header = wl_matching.get_re_tags_with_tokens(self.main, tag_type = 'header')
            RE_TAGS_HEADER = re.compile(tags_header)

            for i, file in enumerate(self.files_to_open):
                self.progress_updated.emit(self.tr('Opening files... ({}/{})').format(i + 1, len_files))

                # Re-decode texts in case encoding settings are manually changed
                file_ext = os.path.splitext(os.path.basename(file['path_orig']))[1].lower()

                if file_ext != '.tmx':
                    file['text'] = get_text_non_tmx(file)
                else:
                    lines = []

                    with open(file['path_orig'], 'r', encoding = file['encoding'], errors = 'replace') as f:
                        soup = bs4.BeautifulSoup(f.read(), 'lxml-xml')

                    for elements_tu in soup.select('tu'):
                        seg_src, seg_tgt = elements_tu.select('seg')

                        if file['tmx_type'] == 'src':
                            lines.append(seg_src.get_text().replace(r'\n', ' ').strip())
                        elif file['tmx_type'] == 'tgt':
                            lines.append(seg_tgt.get_text().replace(r'\n', ' ').strip())

                    file['text'] = '\n'.join(lines)

                # Remove header tags
                with open(file['path'], 'w', encoding = file['encoding']) as f:
                    text = file['text']

                    if file['tagged'] and tags_header:
                        # Use regex here since BeautifulSoup will add tags including <html> and <body> to the text
                        # See: https://www.crummy.com/software/BeautifulSoup/bs4/doc/#differences-between-parsers
                        text = RE_TAGS_HEADER.sub('', text)

                    f.write(text)

                # Process texts
                if self.file_type == 'observed':
                    file['text'] = wl_texts.Wl_Text(self.main, file)
                elif self.file_type == 'ref':
                    file['text'] = wl_texts.Wl_Text_Ref(self.main, file)

                new_files.append(file)
        except Exception:
            err_msg = traceback.format_exc()

        self.progress_updated.emit(self.tr('Updating table...'))
        self.worker_done.emit(err_msg, new_files)
